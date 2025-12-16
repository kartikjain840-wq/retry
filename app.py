import streamlit as st
import json
import time
import sqlite3
import pdfplumber
import docx
from pptx import Presentation
import pandas as pd
import numpy as np
from io import BytesIO
from openai import OpenAI, RateLimitError
from sentence_transformers import SentenceTransformer
import faiss

# ---------------- PAGE CONFIG ----------------
st.set_page_config(page_title="AI Document Intelligence Dashboard", layout="wide")
client = OpenAI(api_key=st.secrets.get("OPENAI_API_KEY"))

# ---------------- DATABASE ----------------
conn = sqlite3.connect("documents.db", check_same_thread=False)
cursor = conn.cursor()

cursor.execute("""
CREATE TABLE IF NOT EXISTS documents (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    filename TEXT,
    objective TEXT,
    tools TEXT,
    results TEXT,
    industry TEXT,
    region TEXT,
    client_type TEXT
)
""")
conn.commit()

# ---------------- SAFE JSON ----------------
def safe_json_list(value):
    if not value:
        return []
    try:
        parsed = json.loads(value)
        return parsed if isinstance(parsed, list) else [str(parsed)]
    except Exception:
        return [str(value)]

# ---------------- BULLET NORMALIZATION ----------------
def normalize_bullets(value):
    bullets = []
    for item in safe_json_list(value):
        parts = [p.strip() for p in str(item).split(",") if p.strip()]
        bullets.extend(parts)
    return bullets

def bullets_to_html(value):
    return "<br>".join([f"• {b}" for b in normalize_bullets(value)])

# ---------------- TEXT EXTRACTION ----------------
def extract_text(file, filetype):
    text = ""
    if filetype == "pdf":
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
    elif filetype == "docx":
        doc = docx.Document(file)
        text = "\n".join(p.text for p in doc.paragraphs)
    elif filetype == "pptx":
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text[:12000]

# ---------------- AI EXTRACTION ----------------
def ai_extract(text):
    prompt = f"""
Return STRICT JSON.
Objective, Tools, Results as bullet lists.

{{
  "objective": [],
  "tools": [],
  "results": [],
  "industry": "",
  "region": "",
  "client_type": ""
}}

TEXT:
{text}
"""
    for _ in range(3):
        try:
            r = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                temperature=0
            )
            return json.loads(r.choices[0].message.content)
        except RateLimitError:
            time.sleep(5)

    return {
        "objective": ["AI extraction failed"],
        "tools": ["AI extraction failed"],
        "results": ["AI extraction failed"],
        "industry": "Unknown",
        "region": "Unknown",
        "client_type": "Unknown"
    }

# ---------------- CONSULTANT'S BRAIN (LAZY RAG) ----------------
@st.cache_resource
def get_rag_components():
    model = SentenceTransformer("all-MiniLM-L6-v2")
    index = faiss.IndexFlatL2(384)
    text_map = []
    return model, index, text_map

# ---------------- UI ----------------
st.title("📊 AI Document Intelligence Dashboard")

search_query = st.text_input("🔍 Search across case studies")

tab1, tab2 = st.tabs(["📂 Upload & Process", "📈 Dashboard"])

# ---------------- TAB 1 ----------------
with tab1:
    uploaded_files = st.file_uploader(
        "Upload PDF / DOCX / PPTX",
        type=["pdf", "docx", "pptx"],
        accept_multiple_files=True
    )

    if uploaded_files and st.button("🚀 Process Documents"):
        embedding_model, faiss_index, faiss_text_map = get_rag_components()

        for file in uploaded_files:
            filename = file.name
            filetype = filename.split(".")[-1].lower()
            text = extract_text(file, filetype)

            ai_data = ai_extract(text)

            cursor.execute("""
                INSERT INTO documents
                (filename, objective, tools, results, industry, region, client_type)
                VALUES (?, ?, ?, ?, ?, ?, ?)
            """, (
                filename,
                json.dumps(ai_data["objective"]),
                json.dumps(ai_data["tools"]),
                json.dumps(ai_data["results"]),
                ai_data["industry"],
                ai_data["region"],
                ai_data["client_type"]
            ))
            conn.commit()

            for chunk in text.split("\n"):
                if len(chunk.strip()) > 60:
                    vec = embedding_model.encode(chunk)
                    faiss_index.add(np.array([vec]).astype("float32"))
                    faiss_text_map.append(chunk)

        st.success("✅ Documents processed and indexed")

# ---------------- TAB 2 ----------------
with tab2:
    df = pd.read_sql("""
        SELECT *
        FROM documents
        GROUP BY filename
        ORDER BY id DESC
    """, conn)

    if df.empty:
        st.info("No documents processed yet.")
    else:
        # -------- INSIGHT CARDS --------
        st.markdown("## 📊 Executive Snapshot")
        c1, c2, c3 = st.columns(3)
        c1.metric("📚 Knowledge Base", f"{len(df)} Case Studies")
        c2.metric("🏭 Top Industry", df["industry"].mode()[0])
        c3.metric(
            "🛠️ Most Used Tool",
            pd.Series(sum([safe_json_list(t) for t in df["tools"]], [])).mode()[0]
        )
        st.markdown("---")

        # -------- FILTERS --------
        f1, f2 = st.columns(2)
        with f1:
            regions = st.multiselect("Filter by Region", sorted(df["region"].unique()))
        with f2:
            industries = st.multiselect("Filter by Industry", sorted(df["industry"].unique()))

        if regions:
            df = df[df["region"].isin(regions)]
        if industries:
            df = df[df["industry"].isin(industries)]

        if search_query:
            df = df[df.apply(
                lambda r: search_query.lower() in " ".join(r.astype(str)).lower(),
                axis=1
            )]

        df["Objectives"] = df["objective"].apply(bullets_to_html)
        df["Results"] = df["results"].apply(bullets_to_html)

        table_df = df[[
            "filename",
            "Objectives",
            "Results",
            "industry",
            "region",
            "client_type"
        ]]

        # -------- DOWNLOAD AS XLSX (ADDITION ONLY) --------
        output = BytesIO()
        export_df = table_df.copy()
        export_df["Objectives"] = export_df["Objectives"].str.replace("<br>", "\n", regex=False)
        export_df["Results"] = export_df["Results"].str.replace("<br>", "\n", regex=False)
        export_df.to_excel(output, index=False, sheet_name="Dashboard")
        output.seek(0)

        st.download_button(
            label="⬇️ Download Dashboard (XLSX)",
            data=output,
            file_name="dashboard_export.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

        # -------- CENTER-ALIGNED TABLE (HTML) --------
        st.markdown(
            """
            <style>
            table {
                width: 100%;
                border-collapse: collapse;
            }
            th, td {
                text-align: center !important;
                vertical-align: middle !important;
                padding: 10px;
            }
            </style>
            """,
            unsafe_allow_html=True
        )

        st.markdown(
            table_df.to_html(escape=False, index=False),
            unsafe_allow_html=True
        )

        # -------- CONSULTANT'S BRAIN --------
        st.markdown("## 🧠 Consultant’s Brain (Ask Like a Partner)")
        user_q = st.text_input(
            "Ask a strategic question",
            placeholder="How did we optimize inventory for manufacturing clients?"
        )

        if user_q:
            embedding_model, faiss_index, faiss_text_map = get_rag_components()

            if faiss_index.ntotal == 0:
                st.warning("No documents indexed yet. Please process documents first.")
            else:
                q_vec = embedding_model.encode(user_q)
                D, I = faiss_index.search(np.array([q_vec]).astype("float32"), k=5)
                context = "\n".join([faiss_text_map[i] for i in I[0]])

                rag_prompt = f"""
You are a senior management consulting partner.
Synthesize an insight-driven answer using the evidence below.

QUESTION:
{user_q}

EVIDENCE:
{context}
"""
                resp = client.chat.completions.create(
                    model="gpt-3.5-turbo",
                    messages=[{"role": "user", "content": rag_prompt}],
                    temperature=0.2
                )

                st.markdown("### 🔍 Partner Synthesis")
                st.write(resp.choices[0].message.content)
